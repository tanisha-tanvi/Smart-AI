import os
import sys
import threading
import platform
import subprocess
from datetime import datetime

from flask import Flask, request, jsonify, render_template
from flask_socketio import SocketIO, emit
from dotenv import load_dotenv
import google.generativeai as genai
from werkzeug.utils import secure_filename

# --- DOCUMENT & TRANSLATION IMPORTS ---
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from deep_translator import GoogleTranslator

# --- INTERNAL SERVICE IMPORTS ---
from services import youtube_service, rag_service, file_creator
from services.metrics_service import tracker
from services.api_manager import api_manager

# --- INITIALIZATION & CONFIG ---
load_dotenv()
WORKSPACE_DIR = os.path.abspath("./workspace")

if not os.path.exists(WORKSPACE_DIR):
    os.makedirs(WORKSPACE_DIR)

# Global configuration for supported file types
SUPPORTED_EXTENSIONS = (
    '.txt', '.py', '.cpp', '.c', '.java', '.js', '.html', '.css', '.md', 
    '.pdf', '.docx', '.pptx', '.jpg', '.jpeg', '.png', '.mp4', '.avi', '.mov'
)

# Initialize Language dictionary for UI
try:
    _supported = GoogleTranslator(source='auto', target='en').get_supported_languages(as_dict=True)
    LANGUAGES = {code: name.title() for name, code in _supported.items()}
except Exception:
    LANGUAGES = {'en': 'English', 'hi': 'Hindi', 'es': 'Spanish', 'fr': 'French'}

# --- DATA MODELS ---
class UserSession:
    """Maintains the state for each active socket connection."""
    def __init__(self):
        self.content = ""
        self.read_pos = 0
        self.lang = 'en'
        self.vector_store = None
        self.current_file = None

user_states = {} 
running_processes = {} 

app = Flask(__name__)
socketio = SocketIO(app, cors_allowed_origins="*", async_mode='threading')

# --- GEMINI TOOL DEFINITIONS ---
tools_schema = [
    {
        "function_declarations": [
            {
                "name": "manage_file",
                "description": "Opens, reads, summarizes, or TRANSCRIBES a local file in the workspace.",
                "parameters": {
                    "type": "OBJECT",
                    "properties": {
                        "filename": {"type": "STRING"},
                        "intent": {"type": "STRING", "enum": ["open", "read", "summarize", "execute", "transcribe"]}
                    },
                    "required": ["filename", "intent"]
                }
            },
            {
                "name": "create_new_file",
                "description": "Creates a new file (docx, py, java, txt) with AI-generated content in the workspace.",
                "parameters": {
                    "type": "OBJECT",
                    "properties": {
                        "filename": {"type": "STRING", "description": "Name including extension, e.g., 'sort.py' or 'essay.docx'"},
                        "prompt": {"type": "STRING", "description": "Description of what content should be inside the file"}
                    },
                    "required": ["filename", "prompt"]
                }
            },
            {
                "name": "youtube_assistant",
                "description": "Plays or transcribes YouTube videos (external links).",
                "parameters": {
                    "type": "OBJECT",
                    "properties": {
                        "query": {"type": "STRING"},
                        "action": {"type": "STRING", "enum": ["play", "transcribe"]}
                    },
                    "required": ["query", "action"]
                }
            },
            {
                "name": "ask_document",
                "description": "Asks a specific question about the currently opened/indexed document.",
                "parameters": {
                    "type": "OBJECT",
                    "properties": {"question": {"type": "STRING"}},
                    "required": ["question"]
                }
            }
        ]
    }
]

def get_rotated_model(use_tools=True):
    """Initializes Gemini with a rotated API key from the pool."""
    api_manager.rotate_key()
    return genai.GenerativeModel(
        'gemini-2.5-flash', 
        tools=tools_schema if use_tools else None
    )

# --- SECURITY & UTILITY HELPERS ---

def get_safe_path(filename):
    """Ensures file operations remain strictly within the workspace directory."""
    safe_name = secure_filename(filename)
    full_path = os.path.abspath(os.path.join(WORKSPACE_DIR, safe_name))
    if not full_path.startswith(WORKSPACE_DIR):
        raise PermissionError("Security Violation")
    return full_path

def extract_text(file_path, sid=None):
    """Parses various document types and triggers RAG indexing automatically."""
    ext = os.path.splitext(file_path)[1].lower()
    text = ""
    try:
        if ext in ['.txt', '.py', '.cpp', '.c', '.java', '.js', '.md', '.html', '.css']:
            with open(file_path, 'r', errors='replace') as f:
                text = f.read()
        elif ext == '.pdf':
            doc = fitz.open(file_path)
            text = "\n".join(page.get_text() for page in doc)
        elif ext == '.docx':
            text = "\n".join([p.text for p in Document(file_path).paragraphs])
        elif ext == '.pptx':
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"): text += shape.text + "\n"
        
        # Build RAG Index in background if text is found
        if text and sid and sid in user_states:
            socketio.emit('terminal_output', {'text': f"[⚙️ Indexing {os.path.basename(file_path)}...]\n"}, room=sid)
            v_store = rag_service.build_rag_index(text)
            if v_store:
                user_states[sid].vector_store = v_store
                user_states[sid].current_file = os.path.basename(file_path)
            else:
                return "ERROR: RAG Indexing failed (Check Embedding API Quota)."
            
        return text
    except Exception as e:
        return f"ERROR: Extraction failed: {str(e)}"

# --- PRIMARY BUSINESS LOGIC ---

def handle_file_intent(filename, intent, sid):
    """Core logic for handling local file requests (Interception layer)."""
    try:
        path = get_safe_path(filename)
    except Exception as e:
        return {"status": "error", "message": str(e)}

    target_lang_name = LANGUAGES.get(user_states[sid].lang, "English")

    if intent == "open":
        try:
            sys_name = platform.system()
            if sys_name == 'Windows': os.startfile(path)
            else: subprocess.call(('open' if sys_name == 'Darwin' else 'xdg-open', path))
            return {"status": "success", "message": f"Successfully opened {filename} locally."}
        except Exception as e:
            return {"status": "error", "message": f"System Open Error: {e}"}

    if intent == "read":
        content = extract_text(path, sid=sid)
        if not content or content.startswith("ERROR:"):
            return {"status": "error", "message": content or "No text found."}
        user_states[sid].content = content
        user_states[sid].read_pos = 0
        return {"status": "success", "action": "start_read", "message": f"Now reading: {filename}"}

    if intent == "summarize":
        try:
            ext = os.path.splitext(filename)[1].lower()
            socketio.emit('terminal_output', {'text': f"[🧠 AI Analyzing {filename} in {target_lang_name}...]\n"}, room=sid)
            
            if ext in ['.docx', '.pptx', '.pdf', '.txt', '.md']:
                text_content = extract_text(path, sid=sid)
                if text_content.startswith("ERROR:"): return {"status": "error", "message": text_content}
                
                model = get_rotated_model(use_tools=False)
                response = api_manager.execute_with_retry(
                    model.generate_content,
                    f"Provide a detailed summary IN {target_lang_name} for the following text:\n\n{text_content}"
                )
                summary = response.text
            else:
                summary = rag_service.summarize_multimodal(path, target_lang=target_lang_name)
                
            user_states[sid].content = summary
            user_states[sid].read_pos = 0
            return {"status": "success", "action": "start_read", "message": summary}
        except Exception as e:
            return {"status": "error", "message": f"AI Summarization Failed: {str(e)}"}

    if intent == "execute":
        docker_service.execute_code_interactive(path, sid, socketio, running_processes)
        return {"status": "success", "message": f"Initiated execution of {filename}."}

    if intent == "transcribe":
        socketio.emit('terminal_output', {'text': f"[🎙️ Local Transcription: Extracting audio from {filename}...]\n"}, room=sid)
        start_time = datetime.now()
        res = youtube_service.transcribe_video_file(path, target_lang=user_states[sid].lang)
        duration = (datetime.now() - start_time).total_seconds()
        
        if res['status'] == 'success':
            tracker.record_transcription(len(res['text_content']), max(0.1, duration))
            user_states[sid].content = res['text_content']
            user_states[sid].read_pos = 0
            return {"status": "success", "action": "start_read", "message": res['text_content']}
        return res

    return {"status": "error", "message": f"Intent '{intent}' not implemented."}

# --- FLASK WEB ROUTES ---

@app.route('/')
def home():
    opts = "".join([f'<option value="{c}" {"selected" if c=="en" else ""}>{n}</option>' 
                   for c, n in LANGUAGES.items()])
    return render_template('index.html', language_options=opts, default_lang='en')

@app.route('/command', methods=['POST'])
def process_command():
    try:
        data = request.json
        sid, query = data.get('sid'), data.get('question', '').strip()
        if sid not in user_states: user_states[sid] = UserSession()
        if 'lang' in data: user_states[sid].lang = data['lang']
        target_lang_name = LANGUAGES.get(user_states[sid].lang, "English")

        parts = query.lower().split(' ', 1)
        action = parts[0]

        # --- MANDATORY INTERCEPTOR ---
        if action in ['open', 'read', 'run', 'execute', 'summarize', 'transcribe'] and len(parts) > 1:
            tracker.record_routing(True)
            return jsonify(handle_file_intent(parts[1].strip(), action if action != "run" else "execute", sid))

        # AI Orchestration with Retry logic
        model = get_rotated_model()
        try:
            response = api_manager.execute_with_retry(
                model.generate_content,
                f"Communication Language: {target_lang_name}. Task/Query: {query}"
            )
        except Exception:
            return jsonify({"status": "error", "message": "Service capacity reached. Retrying with fresh keys..."})
        
        is_fc = bool(response.candidates[0].content.parts[0].function_call)
        tracker.record_routing(is_fc)

        if is_fc:
            fc = response.candidates[0].content.parts[0].function_call
            args = dict(fc.args)
            
            if fc.name == "manage_file":
                return jsonify(handle_file_intent(args['filename'], args.get('intent', 'open'), sid))
            
            if fc.name == "create_new_file":
                socketio.emit('terminal_output', {'text': f"[📄 Generating '{args['filename']}'...]\n"}, room=sid)
                res = file_creator.create_workspace_file(args['filename'], args['prompt'], target_lang=target_lang_name)
                return jsonify(res)

            if fc.name == "ask_document":
                if not user_states[sid].vector_store:
                    return jsonify({"status": "error", "message": "No document indexed. Please 'read' a file first."})
                rag_query = f"ANSWER IN {target_lang_name}: {args.get('question', query)}"
                ans = rag_service.query_rag_document(rag_query, user_states[sid].vector_store)
                return jsonify({"status": "success", "message": ans})

        return jsonify({"status": "info", "message": response.text})

    except Exception as e:
        return jsonify({"status": "error", "message": f"Server Error: {str(e)}"}), 500

@app.route('/list_files')
def get_files():
    try:
        file_list = []
        for root, _, filenames in os.walk(WORKSPACE_DIR):
            for f in filenames:
                if f.lower().endswith(SUPPORTED_EXTENSIONS):
                    rel = os.path.relpath(os.path.join(root, f), WORKSPACE_DIR)
                    file_list.append(rel.replace("\\", "/"))
        return jsonify({"status": "success", "files": file_list})
    except Exception as e:
        return jsonify({"status": "error", "message": str(e)})

@app.route('/read_chunk')
def fetch_text_chunk():
    try:
        sid = request.args.get('sid')
        user = user_states.get(sid)
        if not user or not user.content or user.read_pos >= len(user.content):
            return jsonify({"status": "done"})

        chunk = user.content[user.read_pos : user.read_pos + 600]
        user.read_pos += 600

        if user.lang != 'en':
            try:
                chunk = GoogleTranslator(source='auto', target=user.lang).translate(chunk)
            except: pass

        return jsonify({"status": "reading", "chunk": chunk})
    except:
        return jsonify({"status": "done"})
    
@app.route('/favicon.ico')
def favicon():
    return '', 204

@app.route('/get_metrics')
def get_system_metrics():
    return jsonify(tracker.get_report())

@socketio.on('terminal_input')
def handle_terminal(data):
    sid = request.sid
    if sid in running_processes:
        proc = running_processes[sid]
        if proc.poll() is None:
            proc.stdin.write((data.get('input', '') + "\n").encode())
            proc.stdin.flush()

@socketio.on('kill_process')
def handle_kill():
    sid = request.sid
    if sid in user_states:
        user_states[sid].content = ""
        user_states[sid].read_pos = 0
        
    if sid in running_processes:
        running_processes[sid].terminate()
        del running_processes[sid]
        emit('terminal_output', {'text': "\n[Process Terminated & Speech Silenced]\n"})

@socketio.on('connect')
def handle_connect():
    user_states[request.sid] = UserSession()

if __name__ == '__main__':
    socketio.run(app, host='0.0.0.0', port=5000, debug=True, allow_unsafe_werkzeug=True)
